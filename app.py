
import streamlit as st
# Hide default Streamlit footer
hide_streamlit_style = """
    <style>
        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}
    </style>
"""
#give it a ash background
st.markdown("""
    <style>
        .reportview-container {
            background: #B2BEB5;  
        }
        .main .block-container {
            background: #B2BEB5; 
        }
    </style>
    """,
    unsafe_allow_html=True,
)




st.markdown(hide_streamlit_style, unsafe_allow_html=True)
import base64
from langchain.llms.bedrock import Bedrock
import boto3
from langchain.retrievers import AmazonKendraRetriever
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
import json
from botocore.exceptions import ClientError
import requests
import base64
from requests.exceptions import HTTPError

#for ppt generation
import json
import re
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


if 'custom_details' not in st.session_state:
    st.session_state.custom_details = False
    st.session_state.username = ""
    st.session_state.confluence_url = ""
    st.session_state.confluence_api_token = ""
    st.session_state.space_key = ""
    st.session_state.content = "" # Add this line to save the generated content
    st.session_state.content_ppt = "" # Add this line to save the generated ppt content


# function for encoding image to base64
def get_image_base64(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode()
    
def get_secret():
    secret_name = "kendraRagApp"

    # Create a Secrets Manager client
    session = boto3.session.Session()
    client = session.client(
        service_name='secretsmanager', region_name='ap-south-1'
    )

    try:
        get_secret_value_response = client.get_secret_value(
            SecretId=secret_name
        )
    except ClientError as e:
        raise e

    # Decrypts secret using the associated KMS key.
    secret = get_secret_value_response['SecretString']
    return secret   

# Generating and Publishing new articles to Atlassian Confluence Space
def publish(content, use_custom_details=False, username=None, confluence_url=None, confluence_api_token=None, space_key=None):

    # If not using custom details, get the default details
    if not use_custom_details:
        secrets = json.loads(get_secret())
        username = secrets['username']
        confluence_url = secrets['confluence_space_url']
        confluence_api_token = secrets['confluence_token']
        space_key = secrets['space_key']

    parent_page_id = None #we are publishing to the main content view
    content = content.lstrip("\n")
    title_end_index = content.find("\n")
    new_page_title = content[:title_end_index].strip()
    new_page_content = content[title_end_index:].strip()
    
    auth_str = f"{username}:{confluence_api_token}"

    # Set the API endpoint URL
    url = f"{confluence_url}"
    auth_str_encoded = base64.b64encode(auth_str.encode()).decode()
    # Set the request headers, including the API token or credentials for authentication
    headers = {
        "Authorization": f"Basic {auth_str_encoded}",
        "Content-Type": "application/json"
    }

    # Set the request payload with the new page information
    data = {
        "type": "page",
        "title": new_page_title,
        "space": {"key": space_key},
        "body": {
            "storage": {
                "value": new_page_content,
                "representation": "storage",
            }
        }
    }
    # If the new page should be a child page, specify the parent page ID
    if parent_page_id:
       data["ancestors"]=[{"type": "page", "id": parent_page_id}]


    try:
        # send post request to create the new page
        response = requests.post(url, headers=headers, json=data)
        # If the response was successful, no Exception will be raised
        response.raise_for_status()
    except HTTPError as http_err:
        # If status code is 400, it might be due to duplicate page title
        if response.status_code == 400:
            st.write(f'HTTP error occurred: {http_err}. A page might already exist with the same title.\n Select another query or type in a new query')
        else:
            st.write(f'HTTP error occurred: {http_err}.')
    except Exception as err:
        st.write(f'Other error occurred: {err}.')
    else:
        st.write('Page successfully created.')

#A function that parses the generated text to json format for ppt generation

def parse_format(segment):
        slides = []
        slide = {}
        points = []
        bullet_point_section = False  # For the Bullet Points: format

        for line in segment:
            line = line.strip()

            # Detect start of a new slide or title
            if line.startswith("Slide ") or line.startswith("Title:") or (line and not any(prefix in line for prefix in ["-", "•", "•\t", "Bullet Points:"])):
                # If there's an existing slide, add it to the slides list
                if slide:
                    slide["points"] = points
                    slides.append(slide)
                    slide = {}
                    points = []
                try:
                    slide["subtitle"] = line.split(":")[2].strip() if ":" in line else line
                except:
                    slide["subtitle"] = line.split(":")[1].strip() if ":" in line else line 
            # Detect the subtitle and start a new slide for Format 10
            elif line.startswith("Subtitle:"):
                # If there's an existing slide, add it to the slides list
                if slide:
                    slide["points"] = points
                    slides.append(slide)
                    slide = {}
                    points = []

                slide["subtitle"] = line.split(":")[1].strip()

            # Detect the start of the "Bullet Points:" section
            elif line.startswith("Bullet Points:"):
                bullet_point_section = True

            # Handle bullet points
            elif line.startswith("- ") or line.startswith("• ") or line.startswith("•\t") or bool(re.match(r"^\d+\.", line)) or bullet_point_section:
                point = line.lstrip("-•\t ").split(". ", 1)[-1].strip()
                points.append(point)

            

        # Add the last slide in the segment
        if slide:
            slide["points"] = points
            slides.append(slide)

        # Filter out slides with no points
        slides = [s for s in slides if s.get("points")]

        return {"slides": slides}

def create_ppt_from_json(json_data):
    """
    Create a PowerPoint presentation from parsed JSON data.
    """
    prs = Presentation()

    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    #Access the slide master
    slide_master = prs.slide_master

    #Define footer properties
    footer_text = "Confidential and Proprietary. © 2023 Karaam Analytics. All rights reserved."
    footer_left = Inches(0.5)
    footer_top = prs.slide_height - Inches(1)
    footer_width = prs.slide_width - Inches(1)
    footer_height = Inches(0.5)

    #get capgemini logo
    capgemini_logo = "logo.jpg"

    logo_width = Inches(1.5)  # Adjust as needed
    logo_x_position = prs.slide_width - logo_width - Inches(0.5)  # Adjust the 0.5 inch offset as needed

    logo_height = Inches(0.5)  # Adjust as needed
    logo_y_position = prs.slide_height - logo_height - Inches(1.0)  # Adjust the 0.5 inch offset as needed



    for slide_data in json_data["slides"]:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        #Adjust the width of the content placeholder
        content.width = Inches(12.5)
        content.top = title.top + title.height + Inches(0.5)  # Move content down by 0.5 inches below the title


        title.text = slide_data["subtitle"]
        # Set the title text color to blue
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 85, 183)  # RGB values for blue

        
        for point in slide_data["points"]:
            p = content.text_frame.add_paragraph()
            p.text = point
            p.level = 0  # Bullet point level
            p.space_after = Pt(14) 
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # Add logo to the top right corner
        slide.shapes.add_picture(capgemini_logo, logo_x_position, logo_y_position, width=logo_width)

        # Add footer to the slide
        footer_shape = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
        text_frame = footer_shape.text_frame
        p = text_frame.add_paragraph()
        p.text = footer_text
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        
    #output_filename = "temp_ppt.pptx"
    #prs.save(output_filename)
    
    #return output_filename
    return prs

def qa(query, temperature, topP):
    secrets = json.loads(get_secret())
    kendra_index_id = secrets['kendra_index_id']
    aws_access_key_id = secrets['aws_access_key_id']
    aws_secret_access_key = secrets['aws_secret_access_key']
    BEDROCK_CLIENT = boto3.client('bedrock', region_name='us-east-1', aws_access_key_id=aws_access_key_id,aws_secret_access_key=aws_secret_access_key)
    llm = Bedrock(model_id="amazon.titan-tg1-large", region_name='us-east-1', model_kwargs={"maxTokenCount": 4096, "temperature": temperature, "topP": topP}, client = BEDROCK_CLIENT)
    #top-p is the probability of the most likely token to be sampled, meaning that the model will sample from the most likely tokens with a probability of 0.9, hence the model will be more conservative in its sampling. The range of top-p is between 0 and 1.
    #temperature is a scaling factor for the logits distribution. The range of temperature is between 0 and 1. The lower the temperature, the more conservative the model will be in its sampling. Different temperature values can be used to trade off quality and diversity.
    KENDRA_CLIENT = boto3.client('kendra', region_name='ap-south-1', aws_access_key_id=aws_access_key_id, aws_secret_access_key=aws_secret_access_key)
    

    retriever = AmazonKendraRetriever(index_id=kendra_index_id, region_name='ap-south-1', client=KENDRA_CLIENT)
    
    prompt_template = """
    {context}
    {question} If and only IF you are unable to find the relevant text should you respond with 'Sorry! I can't generate the needed content based on the context provided.'
    """
    
    PROMPT = PromptTemplate(
    template=prompt_template, input_variables=["context", "question"])
    
    chain = RetrievalQA.from_chain_type(
    llm=llm,
    retriever=retriever,
    verbose=True,
    chain_type_kwargs={
    "prompt": PROMPT
    }
    )
    
    return chain(query)

company_logo = get_image_base64("Assets/logo.jpg")
st.markdown(f"""
    <div style="text-align: center;">
        <img src="data:image/jpg;base64,{company_logo}" alt="Karaam Analytics" style="height: 70px; padding-left: 5px;">
        <h4>Retrieval Augmented Generation with Kendra, Amazon Bedrock, and LangChain</h4>      
            </div>
""", unsafe_allow_html=True)

#creating a sidebar with a title and sliders for temperature and top-p
with st.sidebar:
    st.title("Inference Parameters")
    temperature = st.slider("Temperature", min_value=0.0, max_value=1.0, value=0.5, step=0.1)
    topP = st.slider("Top-p", min_value=0.0, max_value=1.0, value=0.5, step=0.1)

#create two radio buttons side by side, one for generating content and the other for generating a ppt
generate_content = st.radio("What do you want to generate?", ("Content", "Powerpoint"))

if generate_content == "Content":

    #creating a text input box for the content query with a default value, the text area should be large enough to fit the query
    input = st.text_area('Enter your query for blog or use default: ', "Please create a well formatted blog titled 'Customer Relations is Pivotal to Business Success' based on the text above. It should have a title. You are allowed to be descriptive and overly long.")
    

    if st.button('Generate'):
        #add a spinner
        with st.spinner('Generating content...'):
            response = qa(input, temperature, topP)
            if response.get("result"):
                st.session_state.content = response["result"] 

    # Outside the 'Generate' button block
    if 'content' in st.session_state:
        st.subheader('Generated Content')
        #display the generated content in a white box
        st.info(st.session_state.content)
        publish_to_confluence = st.radio("Do you want to publish the content to Confluence?", ("No", "Yes"))
        st.session_state.publish_to_confluence = publish_to_confluence

        if publish_to_confluence == "Yes":
            # show the sidebar with the confluence details
            st.sidebar.header("Confluence Space")
            st.session_state.custom_details = st.sidebar.checkbox("Use custom details", st.session_state.custom_details) #session_state.custom_details is a boolean value that is True if the user has selected to use custom details
            if st.session_state.custom_details:
                st.session_state.username = st.sidebar.text_input("Please enter your confluence email:", st.session_state.username)
                st.session_state.confluence_url = st.sidebar.text_input("Please enter your Confluence space url:", st.session_state.confluence_url)
                st.session_state.confluence_api_token = st.sidebar.text_input("Please enter your Confluence API token:", st.session_state.confluence_api_token)
                st.session_state.space_key = st.sidebar.text_input("Please enter your Confluence space key:", st.session_state.space_key)
                #check if the user has entered all the details then enable the publish button
                if st.session_state.username and st.session_state.confluence_url and st.session_state.confluence_api_token and st.session_state.space_key:
                    if st.button("Publish"):
                        publish(st.session_state.content, use_custom_details=True, username=st.session_state.username, confluence_url=st.session_state.confluence_url, confluence_api_token=st.session_state.confluence_api_token, space_key=st.session_state.space_key)
            else:
                st.session_state.username = None
                st.session_state.confluence_url = None
                st.session_state.confluence_api_token = None
                st.session_state.space_key = None
                if st.button("Publish"):
                    publish(st.session_state.content)
        else:
            # if user selects 'No' do nothing
            pass

    else:
        st.subheader('Sorry!')
        st.write("Could not answer the query based on the context available")
elif generate_content == "Powerpoint":
    #create an input for the query with a default value
    input = st.text_area('Enter your query for ppt or use default: ', "Based on the text above, please create a presentation titled 'Customer Relations is Pivotal to Business Success'. For each slide, provide a subtitle followed by bullet points that capture the key ideas. Ensure the content is concise and suitable for a slide format.")
    if st.button('Generate'):
        #add a spinner
        with st.spinner('Generating presentation...'):
            response = qa(input, temperature, topP)
            if response.get("result"):
                st.session_state.content_ppt = response["result"]
            else:
                st.write("Could not generate the ppt based on the available context")
    # Outside the 'Generate' button block
    if 'content_ppt' in st.session_state:
        st.subheader('Generated Presentation')
        #display the generated content in a white box
        st.info(st.session_state.content_ppt)
        #create radio buttons for downloading the ppt
        download_ppt = st.radio("Do you want to generate and download as ppt?", ("No", "Yes"))
        if download_ppt == "Yes":
            #Try to parse the generated text in the correct json format for conversion to ppt
            try:
                print("I entered")
                text=st.session_state.content_ppt
                print("===================================================================")
                #print(text)
                parsed_data = parse_format(text.split("\n"))
                print("I did")
                #create a ppt from the parsed json
                prs = create_ppt_from_json(parsed_data)
                #convert the ppt to bytes
                prs_bytes = BytesIO()
                prs.save(prs_bytes)
                #save to the session state
                st.session_state.ppt_data = prs_bytes.getvalue()

                #create a button to download the ppt
                st.download_button(
                    label="Download ppt",
                    data = st.session_state.ppt_data,
                    #download the ppt
                    file_name="generated_slides.pptx"
                    )
            except Exception as e:
                print(e)
                st.write("Could not parse the content to json or create the ppt. Try again maybe with a different query.")
        else:
            # if user selects 'No' do nothing
            pass



aws_kendra_logo = get_image_base64("Assets/kendra.jpg")
langchain_logo = get_image_base64("Assets/langchain.png")
bedrock_logo = get_image_base64("Assets/bedrock.png")
ec2_logo = get_image_base64("Assets/Amazon_ec2.png")
secrets_manager_logo = get_image_base64("Assets/AWS_secrets.png")
s3_logo = get_image_base64("Assets/s3-bucket.png")

st.markdown(f"""
    <div style="text-align: left;">
        <small>Powered by</small>
        <img src="data:image/jpg;base64,{bedrock_logo}" alt="Amazon Bedrock" style="height: 40px; padding-left: 5px;">
        <img src="data:image/jpg;base64,{aws_kendra_logo}" alt="AWS Kendra" style="height: 40px; padding-left: 5px;">
        <img src="data:image/jpg;base64,{langchain_logo}" alt="Langchain" style="height: 40px; padding-left: 5px;">
        <img src="data:image/jpg;base64,{ec2_logo}" alt="EC2" style="height: 40px; padding-left: 5px;">
        <img src="data:image/jpg;base64,{secrets_manager_logo}" alt="Secrets Manager" style="height: 40px; padding-left: 5px;">
        <img src="data:image/jpg;base64,{s3_logo}" alt="S3" style="height: 40px; padding-left: 5px;">
    </div>
""", unsafe_allow_html=True)
